# -*- coding: utf-8 -*-
"""生成 AI 分享会 PPTX 与 Mermaid 架构图 PNG."""
from __future__ import annotations

import base64
import re
import sys
import textwrap
import urllib.error
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
IMAGES = DOCS / "images" / "ai-sharing"
DIAGRAMS_MD = DOCS / "ai-testing-sharing-diagrams.md"
OUTPUT_PPTX = DOCS / "ai-testing-sharing-session.pptx"

# 幻灯片主题色
BG = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0x4F, 0xC3, 0xF7)
BODY_COLOR = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0x66, 0xBB, 0x6A)

# 页码 -> 插图文件名（None 表示无图）
SLIDE_IMAGES: dict[int, str | None] = {
    4: "diagram-02-toolstack.png",
    5: "diagram-01-pipeline.png",
    8: "diagram-04-human-ai.png",
    10: "diagram-03-api-layers.png",
    12: "diagram-05-jenkins.png",
    13: None,  # 可选 UI 截图，见 _find_ui_screenshot
    16: "diagram-06-case-reuse.png",
    18: "diagram-07-demo-flow.png",
}

SLIDES: list[tuple[str, str, str]] = [
    (
        "测试工程师的 AI 工作流",
        "从需求到自动化回归\n\n基于 Pyautotest 真实项目实践\n\n[你的名字]  |  [部门名称]\n2026 年 [月] 月",
        "开场微笑，停顿 2 秒。今天不讲 AI 概念，只讲手上真在跑的项目。",
    ),
    (
        "我日常在做什么",
        "• CRM 相关测试与自动化（接口 + 业务场景）\n"
        "• 商城 UI 生产巡检（易食包 / ePak）\n"
        "• 移动端稳定性（Monkey / Appium 探索）\n"
        "• 技术栈：Playwright + Pytest + Allure + Jenkins\n\n"
        "项目仓库：Pyautotest",
        "先建立可信度：仓库里真有代码、有 Jenkins、有定时报告。",
    ),
    (
        "测试工作的四个「时间黑洞」",
        "① 测试设计慢 — 需求 → Excel 用例，常常半天起\n"
        "② 自动化起步慢 — 登录、token、定位、框架搭建\n"
        "③ 回归成本高 — 版本一多，手工点不过来\n"
        "④ 协作摩擦大 — 缺截图、缺日志、缺结构化证据",
        "问一句「大家有没有同感」。AI 不是替代测试，是压缩重复劳动。",
    ),
    (
        "我的 AI 工具栈",
        "Cursor IDE — 读整个仓库 · 改多文件 · 跑终端命令\n\n"
        "+ 项目上下文 — 现有代码风格、fixture、报告格式\n\n"
        "+ MCP 扩展（如 TestSprite）— AI 能调用工具\n\n"
        "为什么不是纯 ChatGPT？\n→ 测试是工程问题，需要「上下文」",
        "Cursor 像带代码库记忆的结对同事。右图：三层工具栈。",
    ),
    (
        "一条完整的测试交付链路",
        "产品需求 / Wiki / PRD\n        ↓\n"
        "cases_*.py 结构化数据\n        ↓\n"
        "generate_*_test_artifacts.py\n        ↓\n"
        "Excel 用例 + XMind 脑图\n        ↓\n"
        "pytest 自动化（API / UI / 移动端）\n        ↓\n"
        "Allure · JUnit · 截图 · 邮件\n        ↓\n"
        "Jenkins 定时回归",
        "全场主线。右/下图：总链路架构图。产品记上半段，开发记中间，测试记整条。",
    ),
    (
        "案例 1：需求 → 测试用例资产",
        "输入：CRM 需求文档 / 飞书 Wiki / 评审纪要\n\n"
        "处理：cases_crm313.py 等结构化数据 + AI 扩写边界场景\n\n"
        "输出：Excel 测试用例表 + XMind 测试脑图\n\n"
        "已覆盖 9+ 需求模块",
        "产品同事最该听这段。AI 出初稿 + 人审核，不是直接交最终版。",
    ),
    (
        "结构化用例长什么样",
        "模块    1-外呼手机号解绑\n"
        "用例 ID  OB-003\n"
        "优先级  P0 · 阻塞 · 首轮必测\n"
        "场景    点击解绑 → 二次确认弹窗及文案\n"
        "预期    文案展示实际绑定号码\n\n"
        "规则写清楚 → AI 才能补边界\n"
        "模糊需求 → AI 也模糊",
        "用真实例子比讲概念管用。呼吁产品写清异常规则。",
    ),
    (
        "案例 1：人机分工",
        "AI 擅长                    人必须做\n"
        "────────────────────────────────\n"
        "扩边界场景                  业务规则对不对\n"
        "统一 Excel / 脑图格式       优先级与阻塞判断\n"
        "快速出初稿                  跨模块遗漏检查\n"
        "多需求复用同一套模板        最终 sign-off\n\n"
        "原则：AI 写初稿，人做审核",
        "全场最重要原则之一。右图：人机分工与质量门禁。",
    ),
    (
        "案例 2：用例 → 自动化代码",
        "CRM API 自动化\n  7 个测试文件 · 11 条异常场景\n\n"
        "商城 UI 巡检\n  易食包 / ePak：登录 → 首页 → 商品详情\n\n"
        "移动端\n  Monkey 稳定性 · Appium 随机探索\n\n"
        "技术栈：Playwright + Pytest + Allure",
        "开发同事重点听。接口自动化 ROI 最高。",
    ),
    (
        "CRM API 自动化分层",
        "tests/test_api_*.py — 用例层（含 negative）\n        ↓\n"
        "api/services/*.py — 业务接口封装\n        ↓\n"
        "api/client.py — HTTP 客户端\n        ↓\n"
        "conftest.py fixtures — 登录 token 复用\n        ↓\n"
        "Allure attach — 请求 / 响应 JSON\n\n"
        "异常用例格式固定 → 最适合 AI 批量生成",
        "断言 code=1000 等业务规则必须人审。右图：API 分层。",
    ),
    (
        "不止接口：UI 巡检与移动端",
        "易食包 UI 巡检\n  auth.esbao.com → www.esbao.com 商城\n"
        "  全页滚动 · 热销商品 · 详情页校验 · 截图报告\n\n"
        "ePak 英文商城 — 同类流程\n\n"
        "Monkey — 白屏检测 · 分段截图 · 邮件告警\n\n"
        "难点在稳定性，AI 帮写页面步骤，人调等待与断言",
        "商城巡检 Jenkins 每 30 分钟跑一轮，失败有截图。",
    ),
    (
        "案例 3：接入 Jenkins，变成团队资产",
        "CRM API 回归\n  scripts/run_crm_api_jenkins.bat\n"
        "  支持正/异常用例开关 · JUnit · 邮件\n\n"
        "易食包 UI 巡检\n  Jenkinsfile.esbao-ui · 每 30 分钟触发\n\n"
        "本地能跑不算数\n"
        "CI 能跑、能归档、能通知 → 才算交付",
        "本地和 Jenkins 用同一套脚本。右图：Jenkins 流水线。",
    ),
    (
        "失败了，怎么「说清楚」",
        "reports/ui/esbao/<时间戳>/\n  report.json + 步骤截图\n\n"
        "Allure — 每条 API 用例附 request / response\n\n"
        "reports/test-summary-last.json\n  耗时 · 通过 / 失败统计\n\n"
        "有证据 → 减少测试 / 开发 / 产品三方扯皮",
        "给开发和产品看：打开报告看哪一步、哪个接口、哪张截图。",
    ),
    (
        "MCP 试点：TestSprite",
        "收获\n  自动生成测试计划\n  映射到现有 pytest 用例\n\n"
        "局限\n  默认检测本地端口\n  我们 API 在远程 staging → 曾阻塞\n\n"
        "结论\n  AI 测试工具要和人写代码结合\n  工具输出 ≠ 可以直接上线",
        "诚实讲局限反而加分。TestSprite 报告在 testsprite_tests/ 目录。",
    ),
    (
        "我踩过的四个坑",
        "① 幻觉 — 编造字段/选择器 → 先读现有文件，改完必跑 pytest\n\n"
        "② 一次改太多 — diff 巨大 → 小步迭代，一次一个场景\n\n"
        "③ 盲信工具 — 环境不匹配 → 工具报告 + pytest 交叉验证\n\n"
        "④ 泄密风险 — 账号进 prompt → .env + Jenkins Credentials",
        "每点 15 秒，重点是对策，同事能带走。",
    ),
    (
        "可复制的四条原则",
        "1. AI 写初稿，人 sign-off\n"
        "2. 给足上下文（仓库结构 + 接口文档 + 现有代码）\n"
        "3. 小步提交，立即运行\n"
        "4. 能模板化的先模板化，再规模化\n\n"
        "模板化：cases_*.py → generate 脚本 → Excel / XMind",
        "第四条是精华：9 个需求模块共用同一模式。右图：用例复用模式。",
    ),
    (
        "下周就能试的三件事",
        "测试 — 挑 1 个需求，AI 补 10 条边界用例，审到 5 条\n\n"
        "开发 — 新接口 PR 里带 3 条异常用例，测试 review 后并入回归\n\n"
        "产品 — 需求写清异常规则与「不做什么」\n\n"
        "今天带走一句话：AI 做加速，人做判断",
        "分别看向测试、开发、产品。全场落地页。",
    ),
    (
        "谢谢",
        "接下来：5 分钟 Live Demo\n"
        "  补一条 CRM API 异常用例 → pytest → 看报告\n\n"
        "Q & A 欢迎随时打断\n\n"
        "延伸阅读：docs/ai-testing-sharing-session.md",
        "切到 Cursor 或终端。右图：Demo 闭环序列图。",
    ),
]

DIAGRAM_NAMES = [
    ("diagram-01-pipeline", "图 1：总链路图"),
    ("diagram-02-toolstack", "图 2：工具栈三层图"),
    ("diagram-03-api-layers", "图 3：CRM API 自动化分层"),
    ("diagram-04-human-ai", "图 4：人机分工"),
    ("diagram-05-jenkins", "图 5：Jenkins 回归流水线"),
    ("diagram-06-case-reuse", "图 6：用例生成复用模式"),
    ("diagram-07-demo-flow", "图 7：Demo 闭环"),
]


def _extract_mermaid_blocks(md_text: str) -> list[str]:
    return re.findall(r"```mermaid\s*\n(.*?)```", md_text, re.DOTALL)


def _encode_mermaid_for_ink(source: str) -> str:
    return base64.urlsafe_b64encode(source.encode("utf-8")).decode("ascii")


def export_mermaid_png(source: str, out_path: Path) -> bool:
    encoded = _encode_mermaid_for_ink(source.strip())
    url = f"https://mermaid.ink/img/{encoded}?bgColor=1a1a2e"
    req = urllib.request.Request(url, headers={"User-Agent": "Pyautotest/1.0"})
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            out_path.write_bytes(resp.read())
        return True
    except urllib.error.URLError as exc:
        print(f"  [WARN] mermaid.ink failed for {out_path.name}: {exc}", file=sys.stderr)
        return False


def export_all_diagrams() -> list[Path]:
    IMAGES.mkdir(parents=True, exist_ok=True)
    blocks = _extract_mermaid_blocks(DIAGRAMS_MD.read_text(encoding="utf-8"))
    if len(blocks) < len(DIAGRAM_NAMES):
        print(f"[WARN] expected {len(DIAGRAM_NAMES)} diagrams, found {len(blocks)}", file=sys.stderr)

    exported: list[Path] = []
    for (filename, _label), source in zip(DIAGRAM_NAMES, blocks):
        out = IMAGES / f"{filename}.png"
        print(f"Exporting {out.name} ...")
        if export_mermaid_png(source, out):
            exported.append(out)
            print(f"  OK ({out.stat().st_size // 1024} KB)")
        else:
            print(f"  SKIP")
    return exported


def _find_ui_screenshot() -> Path | None:
    ui_root = ROOT / "reports" / "ui"
    if not ui_root.is_dir():
        return None
    candidates = sorted(ui_root.rglob("*.png"), key=lambda p: p.stat().st_mtime, reverse=True)
    for path in candidates:
        if "failure" not in path.name.lower():
            return path
    return candidates[0] if candidates else None


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text: str, *, font_size: int = 18, bold: bool = False, color: RGBColor = BODY_COLOR):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return box


def _add_title(slide, title: str) -> None:
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(0.35),
        Inches(12.1),
        Inches(0.9),
        title,
        font_size=32,
        bold=True,
        color=TITLE_COLOR,
    )


def _add_image_fit(slide, img_path: Path, left, top, max_width, max_height):
    if not img_path.is_file():
        return
    pic = slide.shapes.add_picture(str(img_path), left, top)
    ratio = min(max_width / pic.width, max_height / pic.height, 1.0)
    if ratio < 1.0:
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = left
    pic.top = top


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ui_shot = _find_ui_screenshot()
    if ui_shot:
        print(f"UI screenshot for slide 13: {ui_shot.relative_to(ROOT)}")
    else:
        print("No UI screenshot found for slide 13 (optional)")

    blank_layout = prs.slide_layouts[6]

    for idx, (title, body, notes) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, BG)
        _add_title(slide, title)

        img_name = SLIDE_IMAGES.get(idx)
        has_side_image = img_name and (IMAGES / img_name).is_file()

        if has_side_image:
            body_width = Inches(5.8)
            _add_textbox(slide, Inches(0.6), Inches(1.35), body_width, Inches(5.8), body, font_size=17)
            _add_image_fit(
                slide,
                IMAGES / img_name,
                Inches(6.6),
                Inches(1.2),
                Inches(6.4),
                Inches(5.9),
            )
        else:
            _add_textbox(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.8), body, font_size=19)

        if idx == 13 and ui_shot:
            _add_image_fit(slide, ui_shot, Inches(7.2), Inches(3.8), Inches(5.5), Inches(3.0))

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


def main() -> int:
    print("=== Export Mermaid diagrams ===")
    exported = export_all_diagrams()
    print(f"Exported {len(exported)} PNG(s) to {IMAGES.relative_to(ROOT)}")

    print("\n=== Build PowerPoint ===")
    pptx_path = build_pptx()
    print(f"PPTX saved: {pptx_path.relative_to(ROOT)}")
    print(f"Slides: {len(SLIDES)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
